"""PowerPoint (.pptx) parser using python-pptx.

Extracts slide text, speaker notes, and embedded tables.
"""

from __future__ import annotations

import logging
from pathlib import Path

from app.parsers.base import (
    BaseParser,
    ExtractedPage,
    ExtractedTable,
    ParseResult,
)

logger = logging.getLogger(__name__)


class PptxParser(BaseParser):
    """Parse PowerPoint presentations."""

    @property
    def supported_mimes(self) -> set[str]:
        return {
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        }

    @property
    def supported_extensions(self) -> set[str]:
        return {".pptx"}

    async def parse(self, file_path: Path) -> ParseResult:
        from pptx import Presentation

        result = ParseResult()

        try:
            prs = Presentation(str(file_path))
        except Exception as exc:
            logger.error("Failed to open PPTX %s: %s", file_path, exc)
            result.warnings.append(f"Failed to open: {exc}")
            return result

        table_idx = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            text_parts: list[str] = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            text_parts.append(para.text.strip())

                if shape.has_table:
                    table = shape.table
                    headers = [
                        cell.text.strip() for cell in table.rows[0].cells
                    ]
                    rows = [
                        [cell.text.strip() for cell in row.cells]
                        for row in table.rows[1:]
                    ]
                    result.tables.append(
                        ExtractedTable(
                            table_index=table_idx,
                            headers=headers,
                            rows=rows,
                            page_number=slide_num,
                        )
                    )
                    table_idx += 1

            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    text_parts.append(f"[Speaker Notes] {notes}")

            result.pages.append(
                ExtractedPage(
                    page_number=slide_num,
                    text="\n".join(text_parts),
                )
            )

        result.page_count = len(result.pages)
        return result
